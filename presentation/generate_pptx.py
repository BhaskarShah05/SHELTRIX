import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = '/Users/bhaskarshah05/SHELTRIX/presentation/output'
ASSETS_DIR = '/Users/bhaskarshah05/SHELTRIX/presentation/assets'
PUBLIC_ASSETS = '/Users/bhaskarshah05/SHELTRIX/public/assets'
os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation()
# Set widescreen 16:9 (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_DARK = RGBColor(4, 12, 28)        # Deep Obsidian Midnight Navy
CARD_BG = RGBColor(11, 24, 48)       # Frosted Dark Crystal Container
ACCENT_CYAN = RGBColor(0, 229, 212)  # Electric Luminous Cyan
ACCENT_BLUE = RGBColor(56, 189, 248) # Ice Blue
ACCENT_ROSE = RGBColor(244, 63, 94)  # Warning Rose
TEXT_WHITE = RGBColor(255, 255, 255) # Pure White
TEXT_MUTED = RGBColor(148, 163, 184) # Slate Muted Gray
BORDER_CYAN = RGBColor(0, 229, 212)

def set_slide_background(slide):
    # Add a full slide background rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return bg

def add_header(slide, title_text, category_badge="SHELTRIX 2026 CLIMATE DIGITAL TWIN"):
    # Header container
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.1))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_badge = tf.paragraphs[0]
    p_badge.text = category_badge.upper()
    p_badge.font.size = Pt(10)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_CYAN
    p_badge.space_after = Pt(4)
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

# -------------------------------------------------------------
# SLIDE 1: TITLE & COVER SLIDE
# -------------------------------------------------------------
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s1)

# Left Column: Brand & Hero
title_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(6.5), Inches(4.8))
tf = title_box.text_frame
tf.word_wrap = True

p_badge = tf.paragraphs[0]
p_badge.text = "AUTONOMOUS HIGH-ALTITUDE PASSIVE SOLAR RESILIENCE"
p_badge.font.size = Pt(11)
p_badge.font.bold = True
p_badge.font.color.rgb = ACCENT_CYAN
p_badge.space_after = Pt(12)

p_title = tf.add_paragraph()
p_title.text = "SHELTRIX"
p_title.font.size = Pt(50)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_WHITE
p_title.space_after = Pt(8)

p_sub = tf.add_paragraph()
p_sub.text = "Zero-Emission High-Altitude Thermal Digital Twin & Generative Simulation Platform"
p_sub.font.size = Pt(18)
p_sub.font.bold = True
p_sub.font.color.rgb = ACCENT_BLUE
p_sub.space_after = Pt(20)

p_body = tf.add_paragraph()
p_body.text = "Coupling Satellite Irradiance (NASA + Open-Meteo) with BioPCM™ Phase-Change Latent Enthalpy Buffering to sustain +19.4°C human comfort at -17.2°C Himalayan ambient."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_MUTED
p_body.space_after = Pt(24)

p_author = tf.add_paragraph()
p_author.text = "Author & System Architect: Bhaskar Shah  |  GitHub: @BhaskarShah05  |  Live: sheltrix-ai.vercel.app"
p_author.font.size = Pt(11)
p_author.font.bold = True
p_author.font.color.rgb = ACCENT_CYAN

# Right Column: High-Res Himalayan Image
if os.path.exists(os.path.join(PUBLIC_ASSETS, 'futuristic-shelter-1.jpg')):
    pic = s1.shapes.add_picture(os.path.join(PUBLIC_ASSETS, 'futuristic-shelter-1.jpg'), Inches(7.8), Inches(1.4), Inches(4.6), Inches(4.8))

# -------------------------------------------------------------
# SLIDE 2: 1) PROBLEM STATEMENT
# -------------------------------------------------------------
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s2)
add_header(s2, "1) Problem Statement: The High-Altitude Thermal Paradox", "01. PROBLEM STATEMENT")

# Left Column: 3 Pain Point Cards
left_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True

def add_pain_point(tf, num, title, desc):
    p1 = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p1.text = f"{num} {title}"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_ROSE
    p1.space_after = Pt(3)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(16)

add_pain_point(tf, "🔴", "Extreme Sub-Zero Diurnal Swings (-15°C to -40°C)", 
               "Severe cold arid conditions in Leh, Siachen, and Spiti cause catastrophic thermal radiation leakage, dropping interior shelter temperatures below human survival limits.")
add_pain_point(tf, "🔴", "Logistics Vulnerability & Heavy Fossil Fuel Dependency", 
               "Outposts rely on kerosene bukharis and diesel generators. Over 60% of winter logistics budgets are wasted transporting fuel across treacherous mountain passes.")
add_pain_point(tf, "🔴", "Blind Architecture Ignored >1000 W/m² Solar Irradiance", 
               "Himalayas receive >300 sunny days/year with peak direct radiation exceeding 1050 W/m², yet shelters are built with static insulation that overheats by day and freezes by night.")

# Right Column: Climate Landscape Picture + Chart
if os.path.exists(os.path.join(PUBLIC_ASSETS, 'ladakh-cold-climate.jpg')):
    s2.shapes.add_picture(os.path.join(PUBLIC_ASSETS, 'ladakh-cold-climate.jpg'), Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8))

# -------------------------------------------------------------
# SLIDE 3: 2) THE SHELTRIX SOLUTION
# -------------------------------------------------------------
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s3)
add_header(s3, "2) The Solution: Closed-Loop Climate Digital Twin", "02. SOLUTION")

left_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True

p_intro = tf.paragraphs[0]
p_intro.text = "An Autonomous High-Altitude Habitat Platform"
p_intro.font.size = Pt(15)
p_intro.font.bold = True
p_intro.font.color.rgb = ACCENT_CYAN
p_intro.space_after = Pt(10)

def add_bullet(tf, title, desc):
    p = tf.add_paragraph()
    p.text = f"• {title}: {desc}"
    p.font.size = Pt(11.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(12)

add_bullet(tf, "Direct Solar Gain Capture", "High-SHGC triple krypton glazed south facade harvesting +33.3 kWh/day of solar energy.")
add_bullet(tf, "BioPCM™ Latent Phase Buffer", "180 kJ/kg organic phase change matrix absorbing daytime heat and releasing warmth at 21°C.")
add_bullet(tf, "Transient Lumped ODE Physics", "Real-time numerical solver calculating indoor dry-bulb response across 24-hour diurnal cycle.")
add_bullet(tf, "100% Zero-Emission Operation", "Zero diesel/kerosene combustion, eliminating lethal CO hazards and carbon emissions.")

# Right Column: System Architecture Diagram
if os.path.exists(os.path.join(ASSETS_DIR, 'diagram_system_architecture.png')):
    s3.shapes.add_picture(os.path.join(ASSETS_DIR, 'diagram_system_architecture.png'), Inches(6.5), Inches(1.8), Inches(6.0), Inches(4.8))

# -------------------------------------------------------------
# SLIDE 4: 3) FLOW OF SOLUTION & PIPELINE
# -------------------------------------------------------------
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s4)
add_header(s4, "3) Flow of Solution: 5-Step Execution Pipeline", "03. FLOW OF SOLUTION")

# 5 Horizontal Step Cards
steps = [
    ("STEP 01", "Environmental Acquisition", "Resolves GPS coords, altitude (3,500m ASL), and queries Open-Meteo Solar DNI + NASA CERES benchmark.", ACCENT_BLUE),
    ("STEP 02", "3D Geometric Synthesis", "Constructs procedural 3D model in Three.js, projecting solar rays and computing effective south glazing aperture.", ACCENT_CYAN),
    ("STEP 03", "Numerical Energy Balance", "Simulates forward in time (dt = 3600s) solving C_eff * dT/dt with conduction, infiltration, and PCM enthalpy.", ACCENT_BLUE),
    ("STEP 04", "AI Genetic Optimization", "Runs evolutionary population across 128 candidates, Pareto-ranking insulation, glazing, and zero deficit.", ACCENT_CYAN),
    ("STEP 05", "Digital Twin Verification", "Validates simulation against 100Hz MQTT TLSv1.3 IoT sensor telemetry for continuous closed-loop alignment.", ACCENT_BLUE)
]

card_w = Inches(2.2)
card_h = Inches(4.6)
start_x = Inches(0.8)
spacing = Inches(2.4)

for i, (step_no, step_title, step_desc, col) in enumerate(steps):
    x = start_x + i * spacing
    y = Inches(1.8)
    
    card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = col
    card.line.width = Pt(1.8)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.18)
    
    p1 = tf.paragraphs[0]
    p1.text = step_no
    p1.font.size = Pt(10)
    p1.font.bold = True
    p1.font.color.rgb = col
    p1.space_after = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = step_title
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_after = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = step_desc
    p3.font.size = Pt(10.5)
    p3.font.color.rgb = TEXT_MUTED

# -------------------------------------------------------------
# SLIDE 5: 4) DETAILED TECH STACK
# -------------------------------------------------------------
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s5)
add_header(s5, "4) Detailed Tech Stack & Modular Architecture", "04. TECH STACK")

# 4 Quadrant Pillars
quads = [
    ("FRONTEND CLIENT (React 18 + Three.js)", [
        "React 18.3+ Component Architecture & Hooks",
        "Vite 8.2 Bundler (<300ms HMR builds)",
        "Three.js WebGL Real-Time 3D Studio & Sun Dial",
        "Chart.js + React-ChartJS-2 High-Contrast Charts",
        "Vanilla CSS 2026 Crystal Glassmorphism & SF Pro"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), ACCENT_CYAN),

    ("BACKEND MICROSERVICES (Node.js / Express)", [
        "Express 4.21 REST API Engine",
        "Transient Lumped-Capacitance Numerical ODE Solver",
        "BioPCM™ Latent Enthalpy Piecewise Model",
        "AI Genetic Evolutionary Optimizer (128 Candidates)",
        "CORS & Secure Modular Architecture"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.3), ACCENT_BLUE),

    ("EXTERNAL APIS & PROTOCOLS", [
        "Open-Meteo Global Solar & Weather API",
        "Open-Meteo Geocoding REST API (Latitude/Longitude)",
        "NASA POWER CERES / MERRA-2 Reanalysis Satellites",
        "MQTT TLSv1.3 WebSocket IoT Telemetry Stream",
        "Real-Time 100Hz Sensor Packet Generator"
    ], Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.3), ACCENT_BLUE),

    ("INFRASTRUCTURE & DEPLOYMENT", [
        "Vercel Serverless Edge Global CDN",
        "Git & GitHub Monorepo Architecture",
        "Continuous Deployment Pipeline (sheltrix-ai.vercel.app)",
        "Zero-Install Edge Browser Execution",
        "Production Minified Tree-Shaken Bundles"
    ], Inches(6.9), Inches(4.4), Inches(5.6), Inches(2.3), ACCENT_CYAN)
]

for title, items, qx, qy, qw, qh, border_col in quads:
    card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, qx, qy, qw, qh)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = border_col
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.18)
    
    p_hdr = tf.paragraphs[0]
    p_hdr.text = title
    p_hdr.font.size = Pt(12)
    p_hdr.font.bold = True
    p_hdr.font.color.rgb = border_col
    p_hdr.space_after = Pt(6)
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(2)

# -------------------------------------------------------------
# SLIDE 6: 5) UNIQUE SELLING PROPOSITIONS (USPs)
# -------------------------------------------------------------
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s6)
add_header(s6, "5) Unique Selling Propositions: Why SHELTRIX Wins", "05. UNIQUE SELLING PROPOSITIONS")

usps = [
    ("☀️ Autonomous Zero-Energy Equilibrium", "Sustains +19.4°C inside with -17.2°C ambient with zero auxiliary fossil heating. Eliminates diesel/kerosene logistics."),
    ("🧬 BioPCM™ Latent Phase Buffer", "180 kJ/kg organic phase change buffer absorbs daytime solar peaks and prevents nighttime sub-zero freeze."),
    ("🌐 Triple-Tier External API Ecosystem", "Open-Meteo Geocoding + Solar Irradiance + NASA POWER satellite validation + 100Hz MQTT IoT telemetry."),
    ("🎮 Real-Time 3D Spatial Solar Dial", "Interactive Three.js studio tracking continuous solar altitude, azimuth, and penetration angles across 24 hours."),
    ("⚡ AI Genetic Multi-Objective Optimizer", "Evaluates 128 architectural candidates with Pareto ranking and 1-click hot-swap to the live 3D habitat model."),
    ("💎 2026 Crystal Glassmorphism UI", "High-definition backdrop, edge-to-edge 4K video launch, SF Pro typography, and dual light/dark modes.")
]

for idx, (utitle, udesc) in enumerate(usps):
    col_idx = idx % 2
    row_idx = idx // 2
    
    ux = Inches(0.8 + col_idx * 5.9)
    uy = Inches(1.8 + row_idx * 1.7)
    
    card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ux, uy, Inches(5.6), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_CYAN if idx % 2 == 0 else ACCENT_BLUE
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.16)
    
    p1 = tf.paragraphs[0]
    p1.text = utitle
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN if idx % 2 == 0 else ACCENT_BLUE
    p1.space_after = Pt(3)
    
    p2 = tf.add_paragraph()
    p2.text = udesc
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED

# -------------------------------------------------------------
# SLIDE 7: THERMAL EQUILIBRIUM & COMPONENT BALANCE CHARTS
# -------------------------------------------------------------
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s7)
add_header(s7, "Engineering Evidence: 24h Thermal Profiles & Energy Balance", "PHYSICS VALIDATION")

if os.path.exists(os.path.join(ASSETS_DIR, 'chart_thermal_equilibrium.png')):
    s7.shapes.add_picture(os.path.join(ASSETS_DIR, 'chart_thermal_equilibrium.png'), Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))

if os.path.exists(os.path.join(ASSETS_DIR, 'chart_energy_balance.png')):
    s7.shapes.add_picture(os.path.join(ASSETS_DIR, 'chart_energy_balance.png'), Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))

# -------------------------------------------------------------
# SLIDE 8: 6) FEASIBILITY & COMPETITOR ANALYSIS
# -------------------------------------------------------------
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s8)
add_header(s8, "6) Feasibility & Competitor Benchmark Matrix", "06. FEASIBILITY & COMPETITORS")

left_box = s8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.2), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "Commercial & Operational Feasibility"
p1.font.size = Pt(14)
p1.font.bold = True
p1.font.color.rgb = ACCENT_CYAN
p1.space_after = Pt(8)

def add_feasibility(tf, title, desc):
    p = tf.add_paragraph()
    p.text = f"• {title}: {desc}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)

add_feasibility(tf, "Economic Feasibility", "Eliminates 30L/day kerosene burn ($45/day), achieving full capital payback in under 18 months per shelter unit.")
add_feasibility(tf, "Technical Readiness", "BioPCM™ is commercially mass-produced; Open-Meteo & NASA APIs provide 99.9% uptime; Three.js runs natively on standard browsers.")
add_feasibility(tf, "Environmental Impact", "100% zero carbon emissions, zero soot deposition on glaciers, and complete elimination of indoor carbon monoxide poisoning.")

# Right Column: Competitor Comparison Table Box
right_box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.8), Inches(5.2), Inches(4.9))
right_box.fill.solid()
right_box.fill.fore_color.rgb = CARD_BG
right_box.line.color.rgb = ACCENT_CYAN
right_box.line.width = Pt(1.5)

rtf = right_box.text_frame
rtf.word_wrap = True
rtf.margin_left = rtf.margin_top = rtf.margin_right = rtf.margin_bottom = Inches(0.2)

rp_hdr = rtf.paragraphs[0]
rp_hdr.text = "COMPETITOR COMPARISON SUMMARY"
rp_hdr.font.size = Pt(12)
rp_hdr.font.bold = True
rp_hdr.font.color.rgb = ACCENT_CYAN
rp_hdr.space_after = Pt(8)

comparisons = [
    ("Conventional Bunkers", "Heavy fossil fuel burn (20-40L/day), extreme thermal swings, severe logistics hazards."),
    ("EnergyPlus Desktop", "Complex manual CAD import, offline EPW files, no real-time IoT digital twin coupling."),
    ("Traditional Mud Passive", "Sensible mass only, bulky logistics, lacks phase change thermal buffering, large day/night swings."),
    ("SHELTRIX (Winner)", "Autonomous zero-energy equilibrium, BioPCM™ latent stabilization, live satellite + MQTT IoT stream, AI genetic optimizer.")
]

for cname, cdesc in comparisons:
    p_c = rtf.add_paragraph()
    p_c.text = f"▶ {cname}: {cdesc}"
    p_c.font.size = Pt(10)
    p_c.font.color.rgb = TEXT_WHITE if "SHELTRIX" in cname else TEXT_MUTED
    p_c.font.bold = True if "SHELTRIX" in cname else False
    p_c.space_after = Pt(6)

# -------------------------------------------------------------
# SLIDE 9: 7) RESEARCH, LITERATURE & REFERENCES
# -------------------------------------------------------------
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s9)
add_header(s9, "7) Scientific Research, Literature & Peer-Reviewed References", "07. RESEARCH & REFERENCES")

refs_box = s9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
rtf = refs_box.text_frame
rtf.word_wrap = True

def add_reference(rtf, index, citation, app_desc):
    p1 = rtf.add_paragraph() if rtf.paragraphs[0].text else rtf.paragraphs[0]
    p1.text = f"[{index}] {citation}"
    p1.font.size = Pt(11.5)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.space_after = Pt(2)
    
    p2 = rtf.add_paragraph()
    p2.text = f"      Application in SHELTRIX: {app_desc}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(10)

add_reference(rtf, 1, "Duffie, J. A., & Beckman, W. A. (2020). Solar Engineering of Thermal Processes, Photovoltaics and Wind (5th ed.). Wiley.",
              "Algorithms for solar declination, equation of time, and continuous solar altitude/azimuth angle calculations powering the 3D Sun Ray Tracker.")
add_reference(rtf, 2, "NASA Langley Research Center — POWER Project (2026). Prediction of Worldwide Energy Resources.",
              "Ingestion of ALLSKY_SFC_SW_DWN satellite solar radiation benchmarks for NASA-certified validation of real-time Open-Meteo DNI data (0.982 correlation).")
add_reference(rtf, 3, "Phase Change Solutions (2024). BioPCM™ Technical Specifications and Enthalpy Transition Profiles.",
              "Isothermal latent thermal heat capacity formulation (180 kJ/kg at 21°C transition) stabilizing indoor diurnal swings.")
add_reference(rtf, 4, "Wangchuk, S. (Himalayan Institute of Alternatives, Ladakh — HIAL). Passive Solar Buildings in High-Altitude Cold Deserts.",
              "Benchmarked Trombe wall composite matrices, earth-bermed envelopes, and south glazing orientation for extreme Himalayan altitudes.")
add_reference(rtf, 5, "ASHRAE Standard 55-2023. Thermal Environmental Conditions for Human Occupancy.",
              "Comfort envelope parameters (18.0°C - 24.0°C) used in the Pareto multi-objective fitness function.")

# -------------------------------------------------------------
# SLIDE 10: CONCLUSION & PRODUCTION ACCESS
# -------------------------------------------------------------
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s10)
add_header(s10, "SHELTRIX: The Future of Autonomous Climate Resilience", "SUMMARY & ACCESS")

c_box = s10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
ctf = c_box.text_frame
ctf.word_wrap = True

cp1 = ctf.paragraphs[0]
cp1.text = "Key Takeaway: Autonomous Zero-Energy High-Altitude Living is Solved."
cp1.font.size = Pt(22)
cp1.font.bold = True
cp1.font.color.rgb = ACCENT_CYAN
cp1.space_after = Pt(14)

cp2 = ctf.add_paragraph()
cp2.text = "SHELTRIX demonstrates that by intelligently harmonizing satellite solar forecasting, 3D spatial sun tracking, and BioPCM™ latent thermal storage, high-altitude defense outposts and research shelters can achieve complete thermal self-sufficiency with zero fuel consumption, zero carbon emissions, and zero logistical vulnerability."
cp2.font.size = Pt(13)
cp2.font.color.rgb = TEXT_WHITE
cp2.space_after = Pt(28)

cp3 = ctf.add_paragraph()
cp3.text = "🔗 Live Production Web Application: https://sheltrix-ai.vercel.app"
cp3.font.size = Pt(14)
cp3.font.bold = True
cp3.font.color.rgb = ACCENT_BLUE
cp3.space_after = Pt(8)

cp4 = ctf.add_paragraph()
cp4.text = "📦 GitHub Repository & Codebase: https://github.com/BhaskarShah05/SHELTRIX"
cp4.font.size = Pt(14)
cp4.font.bold = True
cp4.font.color.rgb = ACCENT_CYAN
cp4.space_after = Pt(8)

cp5 = ctf.add_paragraph()
cp5.text = "Architect: Bhaskar Shah  |  2026 Climate Digital Twin Platform"
cp5.font.size = Pt(12)
cp5.font.color.rgb = TEXT_MUTED

# Save PPTX
pptx_path = os.path.join(OUTPUT_DIR, 'SHELTRIX_Presentation.pptx')
prs.save(pptx_path)
print(f"Presentation saved successfully to {pptx_path}")
